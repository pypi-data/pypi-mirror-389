# FastMDAnalysis/src/fastmdanalysis/utils/slideshow.py
from __future__ import annotations

import sys
import time
from pathlib import Path
from typing import Any, Dict, List, Optional, Sequence, Tuple, Union

# Optional deps
try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN
except Exception as e:  # pragma: no cover
    Presentation = None  # type: ignore[assignment]
    Inches = None        # type: ignore[assignment]
    Pt = None            # type: ignore[assignment]
    PP_ALIGN = None      # type: ignore[assignment]
    _pptx_import_error = e
else:
    _pptx_import_error = None

try:
    from PIL import Image  # type: ignore
except Exception:
    Image = None  # type: ignore[assignment]

try:
    import cairosvg  # type: ignore
except Exception:
    cairosvg = None  # type: ignore

IMG_EXTS: Tuple[str, ...] = (".png", ".jpg", ".jpeg", ".tif", ".tiff", ".bmp", ".gif", ".svg")

__all__ = ["gather_figures", "slide_show"]


def _require_pptx() -> None:
    if Presentation is None:  # pragma: no cover
        raise RuntimeError(
            "python-pptx is required to create slide decks. Install: pip install python-pptx"
        ) from _pptx_import_error


def _stem_to_title(stem: str) -> str:
    s = stem.replace("-", "_")
    parts = [p for p in s.split("_") if p]
    if not parts:
        return stem.title()
    head, rest = parts[0], parts[1:]
    head_up = head.upper() if head.isalpha() and len(head) <= 5 else head.title()
    return f"{head_up} ({' '.join(rest)})" if rest else head_up


def _caption_from_path(p: Path) -> str:
    try:
        rel = p.resolve().relative_to(Path.cwd())
    except Exception:
        rel = p.name
    try:
        size_kb = p.stat().st_size / 1024.0
        mtime = time.strftime("%Y-%m-%d %H:%M:%S", time.localtime(p.stat().st_mtime))
    except Exception:
        size_kb, mtime = 0.0, ""
    return f"{rel} — {size_kb:.0f} KB — saved {mtime}"


def _ensure_png(path: Path, tmpdir: Path) -> Tuple[Path, Optional[str]]:
    """Convert SVG->PNG if cairosvg is available; otherwise return original path."""
    if path.suffix.lower() != ".svg":
        return path, None
    if cairosvg is None:
        return path, f"SVG conversion skipped for {path.name} (install 'cairosvg' to convert)."
    out = tmpdir / f"{path.stem}.png"
    try:
        cairosvg.svg2png(url=str(path), write_to=str(out))
        return out, None
    except Exception as e:  # pragma: no cover
        return path, f"SVG->PNG conversion failed for {path.name}: {e}"


def gather_figures(
    search_roots: Sequence[Union[str, Path]],
    since_epoch: Optional[float] = None,
    exts: Tuple[str, ...] = IMG_EXTS,
) -> List[Path]:
    """
    Recursively gather figure paths from given directories/files, filtered by extension
    and optional modification time lower bound. Sorted oldest->newest by mtime.
    """
    out: List[Path] = []
    for root in search_roots:
        p = Path(root)
        if p.is_file():
            if p.suffix.lower() in exts:
                try:
                    if (since_epoch is None) or (p.stat().st_mtime >= since_epoch):
                        out.append(p)
                except FileNotFoundError:
                    pass
            continue
        if p.is_dir():
            for ext in exts:
                for f in p.rglob(f"*{ext}"):
                    try:
                        if (since_epoch is None) or (f.stat().st_mtime >= since_epoch):
                            out.append(f)
                    except FileNotFoundError:
                        pass
    out.sort(key=lambda q: q.stat().st_mtime if q.exists() else 0.0)
    return out


def _fit_within_emu(aspect: Optional[float], max_w: int, max_h: int) -> Tuple[int, int]:
    """
    Compute width/height (EMU) to fit inside max_w x max_h while preserving aspect ratio.
    If aspect is None/invalid, fall back to the full available box.
    """
    if aspect is None or aspect <= 0:
        return max_w, max_h
    # Try width constraint first
    h_from_w = int(max_w / aspect)
    if h_from_w <= max_h:
        return max_w, h_from_w
    # Else height constraint
    w_from_h = int(max_h * aspect)
    return w_from_h, max_h


def _timestamp_tag() -> str:
    """Return ddmmyy.HHMM in 24h format."""
    return time.strftime("%d%m%y.%H%M")


def slide_show(
    images: Union[str, Path, Sequence[Union[str, Path, Dict[str, Any]]]],
    outpath: Optional[Union[str, Path]] = None,
    title: str = "FastMDAnalysis Results",
    subtitle: Optional[str] = None,
) -> Path:
    """
    Create a PowerPoint deck with one slide per image. Ensures image does not overlap the title.

    Parameters
    ----------
    images : dir path | file path | sequence of paths | sequence of dicts
        If dicts are provided, they may have keys {"path", "title"?, "caption"?}.
    outpath : output .pptx path. If None, saves as 'fastmda_slides_<ddmmyy.HHMM>.pptx' in CWD.
    title   : title slide text
    subtitle: optional subtitle for the title slide

    Returns
    -------
    Path to the created .pptx file.
    """
    _require_pptx()

    # Normalize inputs -> list[{"path": Path, "title": str, "caption": str}]
    records: List[Dict[str, Any]] = []
    if isinstance(images, (str, Path)):
        src = Path(images)
        if src.is_dir():
            files = gather_figures([src])
        elif src.exists():
            files = [src]
        else:
            raise FileNotFoundError(f"No such file or directory: {src}")
        for p in files:
            records.append(
                {"path": Path(p), "title": _stem_to_title(Path(p).stem), "caption": _caption_from_path(Path(p))}
            )
    else:
        for item in images:
            if isinstance(item, (str, Path)):
                p = Path(item)
                records.append({"path": p, "title": _stem_to_title(p.stem), "caption": _caption_from_path(p)})
            elif isinstance(item, dict) and "path" in item:
                p = Path(item["path"])
                t = item.get("title") or _stem_to_title(p.stem)     # <-- fixed 'or'
                c = item.get("caption") or _caption_from_path(p)    # <-- fixed 'or'
                records.append({"path": p, "title": t, "caption": c})
            else:
                raise TypeError("images must be a path or a dict with a 'path' key")

    if not records:
        raise ValueError("No images found to include in the slide deck.")

    # Output filename with timestamp if not provided
    if outpath is None:
        stamp = _timestamp_tag()  # ddmmyy.HHMM
        outpath = Path.cwd() / f"fastmda_slides_{stamp}.pptx"
    else:
        outpath = Path(outpath)
        if outpath.is_dir():
            stamp = _timestamp_tag()
            outpath = outpath / f"fastmda_slides_{stamp}.pptx"
        if outpath.suffix.lower() != ".pptx":
            outpath = outpath.with_suffix(".pptx")

    prs = Presentation()

    # Title slide
    title_layout = prs.slide_layouts[0]  # Title + Subtitle
    s0 = prs.slides.add_slide(title_layout)
    s0.shapes.title.text = title
    s0.placeholders[1].text = subtitle or f"{len(records)} figure(s) — generated {time.strftime('%Y-%m-%d %H:%M:%S')}"

    # Temp dir for conversions (in cwd)
    tmpdir = Path.cwd() / ".fastmda_tmp"
    tmpdir.mkdir(exist_ok=True)

    # Geometry constants (EMU)
    slide_w, slide_h = prs.slide_width, prs.slide_height
    left_margin = int(Inches(0.5))
    right_margin = int(Inches(0.5))
    top_margin_min = int(Inches(0.3))       # minimal padding above image area
    title_gap = int(Inches(0.15))           # gap between title placeholder and image
    caption_height = int(Inches(0.6))       # reserved caption box height
    caption_gap = int(Inches(0.15))         # gap between image and caption
    bottom_margin = int(Inches(0.5))

    warn_msgs: List[str] = []

    for rec in records:
        p: Path = Path(rec["path"])
        if not p.exists():
            warn_msgs.append(f"Missing image skipped: {p}")
            continue

        # Convert SVG if possible
        use_img, wmsg = _ensure_png(p, tmpdir)
        if wmsg:
            warn_msgs.append(wmsg)

        # Read aspect ratio if PIL is available
        aspect: Optional[float] = None
        if Image is not None:
            try:
                with Image.open(use_img) as im:
                    w_px, h_px = im.size
                    aspect = (w_px / h_px) if h_px else None
            except Exception:
                pass

        # Choose layout with title area
        layout = prs.slide_layouts[5]  # Title Only
        s = prs.slides.add_slide(layout)

        # Title: ensure presence and get its bottom edge
        title_shape = s.shapes.title
        if title_shape is None:
            # Create a simple title box if layout didn't provide one
            title_shape = s.shapes.add_textbox(
                left=int(Inches(0.5)),
                top=int(Inches(0.2)),
                width=int(slide_w - Inches(1.0)),
                height=int(Inches(0.6)),
            )
        title_shape.text = str(rec.get("title") or _stem_to_title(use_img.stem))
        title_bottom = int(title_shape.top + title_shape.height)

        # Compute content region below the title, with padding and reserved caption space
        content_left = left_margin
        content_right = slide_w - right_margin
        content_width = int(content_right - content_left)

        img_top_min = max(top_margin_min, title_bottom + title_gap)
        content_bottom = slide_h - (bottom_margin + caption_gap + caption_height)
        max_img_height = int(content_bottom - img_top_min)
        if max_img_height < int(Inches(1.0)):
            # If the title is unusually tall, degrade gracefully
            max_img_height = int(slide_h - (title_bottom + title_gap + bottom_margin))

        fit_w, fit_h = _fit_within_emu(aspect, content_width, max_img_height)
        left = int(content_left + (content_width - fit_w) / 2)
        top = int(img_top_min + (max_img_height - fit_h) / 2)

        # Place image
        s.shapes.add_picture(str(use_img), left=left, top=top, width=fit_w, height=fit_h)

        # Caption below image (centered)
        cap_text = str(rec.get("caption") or _caption_from_path(use_img))
        cap_left = int(Inches(0.5))
        cap_top = int(top + fit_h + caption_gap)
        cap_width = int(slide_w - Inches(1.0))
        cap_h = caption_height
        # Clamp caption to bottom area if needed
        cap_top = min(cap_top, int(slide_h - bottom_margin - cap_h))

        cap_box = s.shapes.add_textbox(cap_left, cap_top, cap_width, cap_h)
        tf = cap_box.text_frame
        tf.clear()
        para = tf.paragraphs[0]
        para.text = cap_text
        if Pt is not None:
            para.font.size = Pt(12)
        if PP_ALIGN is not None:
            para.alignment = PP_ALIGN.CENTER

    prs.save(str(outpath))

    # Attempt to remove tmpdir if empty
    try:
        any_files = any(tmpdir.iterdir())
    except Exception:
        any_files = True
    if not any_files:
        try:
            tmpdir.rmdir()
        except Exception:
            pass

    # Emit warnings to stderr
    for m in warn_msgs:
        print(f"[slide_show] WARNING: {m}", file=sys.stderr)

    return outpath

