#!/usr/bin/env python3
"""
PowerPoint Post-Processing Script

This script processes PowerPoint presentations to automatically enable text wrapping
and shrink text on overflow for text boxes that contain text longer than a specified threshold.
"""

import argparse
import os
import sys
from pathlib import Path
from typing import Optional

from pptx import Presentation
from pptx.enum.text import MSO_AUTO_SIZE
from pptx.shapes.base import BaseShape
from pptx.text.text import TextFrame

from .config import Config


class PowerPointPostProcessor:
    """Post-processor for PowerPoint presentations to handle text box auto-fitting."""
    
    def __init__(self, config: Optional[Config] = None, verbose: bool = False):
        """Initialize the post-processor with configuration."""
        self.config = config or Config()
        self.verbose = verbose
        self.enable_autofit = self.config.get_bool('ENABLE_TEXT_AUTOFIT', True)
        self.text_threshold = self.config.get_int('TEXT_LENGTH_THRESHOLD', 10)
        
    def process_presentation(self, input_file: str, output_file: Optional[str] = None) -> str:
        """
        Process a PowerPoint presentation to enable text auto-fitting.
        
        Args:
            input_file: Path to the input PowerPoint file
            output_file: Path to save the processed file (optional)
            
        Returns:
            Path to the processed file
        """
        if not os.path.exists(input_file):
            raise FileNotFoundError(f"Input file not found: {input_file}")
            
        # Use input file as output file if not provided (overwrite original)
        if not output_file:
            output_file = input_file
        
        print(f"Processing PowerPoint file: {input_file}")
        print(f"Text length threshold: {self.text_threshold} characters")
        print(f"Auto-fit enabled: {self.enable_autofit}")
        
        # Load presentation
        presentation = Presentation(input_file)
        
        total_processed = 0
        total_slides = len(presentation.slides)
        
        # Process each slide
        for slide_idx, slide in enumerate(presentation.slides, 1):
            if self.verbose:
                print(f"Processing slide {slide_idx}/{total_slides}...")
            processed_count = self._process_slide(slide)
            total_processed += processed_count
            
            if processed_count > 0 and self.verbose:
                print(f"  → Processed {processed_count} text boxes")
        
        # Save the processed presentation
        presentation.save(output_file)
        
        if self.verbose:
            print(f"\nPost-processing completed!")
            print(f"Total text boxes processed: {total_processed}")
            print(f"Output saved to: {output_file}")
        
        return output_file
    
    def _process_slide(self, slide) -> int:
        """
        Process a single slide to enable text auto-fitting for qualifying text boxes.
        
        Args:
            slide: PowerPoint slide object
            
        Returns:
            Number of text boxes processed
        """
        processed_count = 0
        
        for shape in slide.shapes:
            if self._should_process_shape(shape):
                self._apply_text_autofit(shape)
                processed_count += 1
                
        return processed_count
    
    def _should_process_shape(self, shape: BaseShape) -> bool:
        """
        Determine if a shape should be processed for text auto-fitting.
        
        Args:
            shape: PowerPoint shape object
            
        Returns:
            True if the shape should be processed
        """
        if not self.enable_autofit:
            return False
            
        # Check if shape has text frame
        if not hasattr(shape, 'text_frame') or shape.text_frame is None:
            return False
            
        text_frame = shape.text_frame
        
        # Get all text content from the text frame
        text_content = self._get_text_content(text_frame)
        
        # Check if text length exceeds threshold
        if len(text_content.strip()) <= self.text_threshold:
            return False
            
        return True
    
    def _get_text_content(self, text_frame: TextFrame) -> str:
        """
        Extract all text content from a text frame.
        
        Args:
            text_frame: PowerPoint text frame object
            
        Returns:
            Combined text content
        """
        text_content = ""
        
        for paragraph in text_frame.paragraphs:
            for run in paragraph.runs:
                text_content += run.text
                
        return text_content
    
    def _apply_text_autofit(self, shape: BaseShape) -> None:
        """
        Apply text auto-fitting settings to a shape.
        
        Args:
            shape: PowerPoint shape object
        """
        try:
            text_frame = shape.text_frame
            
            # Enable text wrapping in shape
            text_frame.word_wrap = True
            
            # Enable shrink text on overflow
            text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
            
            # Optional: Set margins to provide better text fitting
            # You can adjust these values as needed
            text_frame.margin_left = text_frame.margin_left or 91440  # 0.1 inch in EMUs
            text_frame.margin_right = text_frame.margin_right or 91440
            text_frame.margin_top = text_frame.margin_top or 45720   # 0.05 inch in EMUs
            text_frame.margin_bottom = text_frame.margin_bottom or 45720
            
        except Exception as e:
            print(f"Warning: Could not apply auto-fit to shape: {e}")


def main():
    """Main function for command-line usage."""
    parser = argparse.ArgumentParser(
        description="Post-process PowerPoint presentations to enable text auto-fitting"
    )
    parser.add_argument(
        "--input-file", "-i",
        required=True,
        help="Path to the input PowerPoint file (.pptx)"
    )
    parser.add_argument(
        "--output-file", "-o",
        help="Path to save the processed file (optional, overwrites input file if not provided)"
    )
    parser.add_argument(
        "--text-threshold", "-t",
        type=int,
        help="Text length threshold for enabling auto-fit (overrides .env setting)"
    )
    parser.add_argument(
        "--disable-autofit",
        action="store_true",
        help="Disable text auto-fitting (for testing purposes)"
    )
    parser.add_argument(
        "--debug",
        action="store_true",
        help="Enable debug output"
    )
    
    args = parser.parse_args()
    
    try:
        # Load configuration
        config = Config()
        
        # Override config with command-line arguments
        if args.text_threshold is not None:
            config.set('TEXT_LENGTH_THRESHOLD', str(args.text_threshold))
        if args.disable_autofit:
            config.set('ENABLE_TEXT_AUTOFIT', 'false')
        if args.debug:
            config.set('DEBUG', 'true')
        
        # Create post-processor
        processor = PowerPointPostProcessor(config)
        
        # Process the presentation
        output_file = processor.process_presentation(args.input_file, args.output_file)
        
        print(f"\n✅ Post-processing completed successfully!")
        print(f"📁 Output file: {output_file}")
        
    except Exception as e:
        print(f"❌ Error during post-processing: {e}")
        if args.debug:
            import traceback
            traceback.print_exc()
        sys.exit(1)


# Alias for backward compatibility
PostProcessor = PowerPointPostProcessor


if __name__ == "__main__":
    main()
