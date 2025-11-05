import collections
import collections.abc
from pptx.enum.shapes import MSO_SHAPE


from timeline.utils.shapes import set_shape_transparency, send_backwards, add_text_box
from timeline.utils.paragraph import add_paragraph, amend_font


from pptx.dml.color import RGBColor
from webcolors import hex_to_rgb


def move_elements_to_right(ppt, sidebar_width=0.12):
    """moves all placeholders and images to the right of the screen, shrinking them while maintaining their original aspect ratio, and ensuring their vertical centers remain in the same position"""
    for slide in ppt.slides:
        for shape in slide.shapes:
            # original dimensions and positions
            original_top = shape.top
            original_left = shape.left
            original_height = shape.height
            original_width = shape.width

            # calculate scale factor for width
            content_space_width = ppt.slide_width * (1 - sidebar_width)
            scale_factor = content_space_width / ppt.slide_width

            # new dimensions
            new_width = original_width * scale_factor
            new_height = original_height * scale_factor

            # new positions
            new_left = ppt.slide_width * sidebar_width + (original_left - original_width / 2) * scale_factor + new_width / 2
            vertical_center_offset = (original_height - new_height) / 2
            new_top = original_top + vertical_center_offset

            # apply new dimensions and positions
            shape.left, shape.width = int(new_left), int(new_width)
            shape.top, shape.height = int(new_top), int(new_height)


def merge_tags(tags: list[str]) -> list[str]:
    """merge the adjacent slides with the same tag"""
    merged_tags = []
    for i, tag in enumerate(tags):
        if i == 0:
            merged_tags.append(tag)
        elif tag != merged_tags[-1]:
            merged_tags.append(tag)
    return merged_tags


def set_sidebar_timeline(
    ppt,
    tags: list[str],
    sidebar_width=0.12,
    sidebar_transparency=50000,
    sidebar_color=RGBColor(*hex_to_rgb("#5A5A5A")),
    sidebar_color_outline=RGBColor(*hex_to_rgb("#FFFFFF")),
    sidebar_item_height=0.06,
    sidebar_init_font_size=16,
    sidebar_item_font="Arial",
    sidebar_item_font_color=RGBColor(255, 255, 255),
    indicator_color=RGBColor(*hex_to_rgb("#111111")),
    indicator_transparency=80000,
):
    assert len(tags) == len(ppt.slides), f"The number of tags: {len(tags)} has to match the number of slides: {len(ppt.slides)}"

    merged_tags = merge_tags(tags)

    # adding the base shapes to each slide
    for slide in ppt.slides:
        # ------------------------ shaping the sidebar itself ------------------------
        sidebar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, ppt.slide_width * sidebar_width, ppt.slide_height)
        sidebar.name = "!!SIDEBAR"
        sidebar.fill.solid()
        sidebar.fill.fore_color.rgb = sidebar_color
        sidebar.line.color.rgb = sidebar_color_outline
        set_shape_transparency(sidebar, sidebar_transparency)

        offset = 0
        for tag in merged_tags:
            text_box = add_text_box(
                slide=slide,
                ppt=ppt,
                left=0,
                top=offset,
                width=sidebar_width,
                height=sidebar_item_height,
            )

            add_paragraph(
                placeholder=text_box,
                text=tag,
                font_size=sidebar_init_font_size,
                font_family=sidebar_item_font,
                font_color=sidebar_item_font_color,
            )

            setattr(text_box, "name", f"!!SIDEBAR_{merged_tags.index(tag)}")

            if tags[ppt.slides.index(slide)] == tag:
                amend_font(
                    placeholder=text_box,
                    font_family=sidebar_item_font,
                    font_size=sidebar_init_font_size + 3,
                    bold=True,
                )

                indicator = slide.shapes.add_shape(
                    MSO_SHAPE.RECTANGLE,
                    left=0,
                    top=offset * ppt.slide_height,
                    width=ppt.slide_width * (sidebar_width + 0.01),
                    height=ppt.slide_height * sidebar_item_height,
                )

                indicator.name = "!!INDICATOR"
                indicator.fill.solid()
                indicator.fill.fore_color.rgb = indicator_color
                indicator.line.color.rgb = sidebar_color_outline
                set_shape_transparency(sidebar, indicator_transparency)
                send_backwards(slide, indicator)

            offset += sidebar_item_height

        send_backwards(slide, sidebar)
