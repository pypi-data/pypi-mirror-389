"""Extract subchunks from PPTX file."""

from pathlib import Path
from typing import Any

from pptx import Presentation

from unichunking.tools import (
    compact_matrix,
    convert_file,
    extract_charts,
)
from unichunking.types import (
    ChunkPosition,
    MatrixTable,
    SubChunk,
)

MSO_SHAPE_TYPE_GROUP = 6


def _intersect(
    shape_1_pos: ChunkPosition,
    shape_2_pos: ChunkPosition,
) -> bool:
    shape_1_bary = (
        (shape_1_pos.x0 + shape_1_pos.x1) / 2,
        (shape_1_pos.y0 + shape_1_pos.y1) / 2,
    )
    shape_2_bary = (
        (shape_2_pos.x0 + shape_2_pos.x1) / 2,
        (shape_2_pos.y0 + shape_2_pos.y1) / 2,
    )
    return (
        (shape_1_pos.x0 <= shape_2_bary[0] <= shape_1_pos.x1)
        and (shape_1_pos.y0 <= shape_2_bary[1] <= shape_1_pos.y1)
    ) or (
        (shape_2_pos.x0 <= shape_1_bary[0] <= shape_2_pos.x1)
        and (shape_2_pos.y0 <= shape_1_bary[1] <= shape_2_pos.y1)
    )


def _create_shape_groups(
    shapes: list[Any],
    dimensions: ChunkPosition,
) -> list[tuple[ChunkPosition, list[tuple[Any, ChunkPosition]]]]:
    formatted_shapes_info: list[
        tuple[ChunkPosition, list[tuple[Any, ChunkPosition]]]
    ] = []
    while shapes:
        shape = shapes.pop()
        position = ChunkPosition(
            (shape.left - dimensions.x0) / (dimensions.x1 - dimensions.x0),
            (shape.top - dimensions.y0) / (dimensions.y1 - dimensions.y0),
            (shape.left + shape.width - dimensions.x0)
            / (dimensions.x1 - dimensions.x0),
            (shape.top + shape.height - dimensions.y0)
            / (dimensions.y1 - dimensions.y0),
        )
        formatted_shapes_info.append((position, [(shape, position)]))

    changed_something = True
    while changed_something:
        changed_something = False
        for i in range(len(formatted_shapes_info)):
            for j in range(len(formatted_shapes_info)):
                if i == j:
                    continue
                if _intersect(formatted_shapes_info[i][0], formatted_shapes_info[j][0]):
                    position = ChunkPosition(
                        x0=min(
                            formatted_shapes_info[i][0].x0,
                            formatted_shapes_info[j][0].x0,
                        ),
                        y0=min(
                            formatted_shapes_info[i][0].y0,
                            formatted_shapes_info[j][0].y0,
                        ),
                        x1=max(
                            formatted_shapes_info[i][0].x1,
                            formatted_shapes_info[j][0].x1,
                        ),
                        y1=max(
                            formatted_shapes_info[i][0].y1,
                            formatted_shapes_info[j][0].y1,
                        ),
                    )
                    shapes_list = formatted_shapes_info[i][1]
                    shapes_list.extend(formatted_shapes_info[j][1])
                    formatted_shapes_info[i] = (position, shapes_list)
                    formatted_shapes_info.pop(j)
                    changed_something = True
                    break
            if changed_something:
                break

    return formatted_shapes_info


def _insert_in_order(
    position: ChunkPosition,
    shape_list: list[Any],
    ordered_shape_groups: list[tuple[ChunkPosition, list[Any]]],
    verticality_bias: float = 0.125,
) -> list[tuple[ChunkPosition, list[Any]]]:
    coord = ((position.x0 + position.x1) * (1 - verticality_bias) / 2) ** 2 + (
        (position.y0 + position.y1) * (1 + verticality_bias) / 2
    ) ** 2
    i = 0
    while i < len(ordered_shape_groups):
        if (
            coord
            < (
                (ordered_shape_groups[i][0].x0 + ordered_shape_groups[i][0].x1)
                * (1 - verticality_bias)
                / 2
            )
            ** 2
            + (
                (ordered_shape_groups[i][0].y0 + ordered_shape_groups[i][0].y1)
                * (1 + verticality_bias)
                / 2
            )
            ** 2
        ):
            return (
                ordered_shape_groups[:i]
                + [(position, shape_list)]
                + ordered_shape_groups[i:]
            )
        i += 1
    return [*ordered_shape_groups, (position, shape_list)]


def _sort_shapes(
    shapes: list[Any],
    dimensions: ChunkPosition,
) -> list[Any]:
    """Sorts the shapes on a slide to emulate a natural reading order.

    Overlapping shapes are grouped, and the groups are sorted by increasing distance from their barycenter to the top left corner of the slide.
    Then the shapes in each group are sorted recursively, the coordinates of the group replacing those of the slide, etc.
    The parameter vertical_bias in insert_in_order gives a stronger weight to the vertical axis than to the horizontal one when computing distances.

    Args:
        shapes: List of shapes to sort.
        dimensions: The dimensions of the plane on which to sort the shapes (the slide, or a shape group).

    Returns:
        A list of ordered shapes.
    """
    if len(shapes) <= 1:
        if shapes and shapes[0].shape_type == MSO_SHAPE_TYPE_GROUP:
            return _sort_shapes(list(shapes[0].shapes), dimensions)
        return shapes

    shape_groups = _create_shape_groups(shapes, dimensions)

    ordered_shape_groups: list[tuple[ChunkPosition, list[Any]]] = []

    if len(shape_groups) == 1:
        shape_group = shape_groups[0]
        shapes_in_group = shape_group[1]
        for shape_info in shapes_in_group:
            shape = shape_info[0]
            shape_position = shape_info[1]
            if shape.shape_type == MSO_SHAPE_TYPE_GROUP:
                ordered_shape_groups = _insert_in_order(
                    position=shape_position,
                    shape_list=_sort_shapes(list(shape.shapes), shape_position),
                    ordered_shape_groups=ordered_shape_groups,
                )
            else:
                ordered_shape_groups = _insert_in_order(
                    position=shape_position,
                    shape_list=[shape],
                    ordered_shape_groups=ordered_shape_groups,
                )

    else:
        for shape_group in shape_groups:
            group_position = shape_group[0]
            shapes_in_group_info = shape_group[1]
            shapes_in_group = [
                shape_in_group_info[0] for shape_in_group_info in shapes_in_group_info
            ]
            ordered_shape_groups = _insert_in_order(
                position=group_position,
                shape_list=_sort_shapes(shapes_in_group, group_position),
                ordered_shape_groups=ordered_shape_groups,
            )

    ordered_shapes: list[Any] = []
    for shape_group in ordered_shape_groups:
        ordered_shapes.extend(shape_group[1])

    return ordered_shapes


def _handle_text_frame(
    shape: Any,
    subchunks: list[SubChunk],
    slide_idx: int,
    position: ChunkPosition,
    filename: str,
) -> list[SubChunk]:
    text_frame = shape.text_frame
    for paragraph in text_frame.paragraphs:
        for run in paragraph.runs:
            if run.text:
                text = run.text
                if run.font and run.font.bold:
                    text = f"**{text}**"
                subchunks.append(
                    SubChunk(
                        subchunk_id=len(subchunks),
                        content=text + " ",
                        page=slide_idx,
                        position=position,
                        file_name=filename,
                    ),
                )
    return subchunks


def _handle_charts_and_tables(
    shape: Any,
    subchunks: list[SubChunk],
    tables: list[MatrixTable],
    slide_idx: int,
    position: ChunkPosition,
    filename: str,
    file_path: Path,
) -> tuple[list[SubChunk], list[MatrixTable]]:
    if shape.has_chart:
        for plot in extract_charts(
            file_path,
            shape._element.xml,  # noqa: SLF001
            "pptx",
            slide_idx + 1,
        ):
            subchunks.append(
                SubChunk(
                    subchunk_id=len(subchunks),
                    content=f"{len(tables)}",
                    page=slide_idx,
                    position=position,
                    file_name=filename,
                    content_type="table",
                ),
            )
            tables.append(plot)

    if shape.has_table:
        table = shape.table

        matrix_table = [
            [str(table.cell(i, j).text) for j in range(len(table.columns))]
            for i in range(len(table.rows))
        ]
        matrix_table = compact_matrix(table=MatrixTable("", matrix_table))
        if matrix_table.abort():
            subchunks.append(
                SubChunk(
                    subchunk_id=len(subchunks),
                    content=matrix_table.read_as_str(),
                    page=slide_idx,
                    position=position,
                    file_name=filename,
                ),
            )
        else:
            subchunks.append(
                SubChunk(
                    subchunk_id=len(subchunks),
                    content=f"{len(tables)}",
                    page=slide_idx,
                    position=position,
                    file_name=filename,
                    content_type="table",
                ),
            )
            tables.append(matrix_table)

    return subchunks, tables


async def extract_subchunks_pptx(
    path: Path,
    extension: str,
) -> tuple[list[SubChunk], list[MatrixTable], list[str], Path, bool]:
    """Filetype-specific function : extracts subchunks from a PPTX file.

    For PPT & ODP : local file converted to PPTX, subchunks extracted from PPTX.
    Keeps a numbered slides structure.
    (ex : "Slide 1 : ... ; End of slide 1").

    Args:
        path: Path to the local file.
        extension: File extension.
        function: Function to handle images.

    Returns:
        A tuple containing three lists and a Path:
        - List of subchunks, containing actual text or markers pointing to table/chart/image objects.
        - List of table/chart objects, of class MatrixTable.
        - List of image objects.
        - Path to the processed PPTX file, different from the initial path if a conversion occured.
    """
    if extension != "pptx":
        new_path = await convert_file(path, "pptx")
        if new_path is not None:
            path = new_path
        else:
            return [], [], [], Path(), False

    prs = Presentation(str(path))

    subchunks: list[SubChunk] = []
    tables: list[MatrixTable] = []

    for slide_idx in range(len(prs.slides)):
        slide = prs.slides[slide_idx]
        subchunks.append(
            SubChunk(
                subchunk_id=len(subchunks),
                content=f"** Slide {slide_idx + 1} : ",
                page=slide_idx,
                position=ChunkPosition(0, 0, 0, 0),
                file_name=path.name,
            ),
        )

        slide_width = float(prs.slide_width or 1)
        slide_height = float(prs.slide_height or 1)
        ordered_shapes = _sort_shapes(
            shapes=list(slide.shapes),
            dimensions=ChunkPosition(0, 0, slide_width, slide_height),
        )

        for shape in ordered_shapes:
            position = ChunkPosition(
                shape.left / slide_width,
                shape.top / slide_height,
                (shape.left + shape.width) / slide_width,
                (shape.top + shape.height) / slide_height,
            )

            subchunks, tables = _handle_charts_and_tables(
                shape=shape,
                subchunks=subchunks,
                tables=tables,
                slide_idx=slide_idx,
                position=position,
                filename=path.name,
                file_path=path,
            )

            if shape.has_text_frame:
                subchunks = _handle_text_frame(
                    shape=shape,
                    subchunks=subchunks,
                    slide_idx=slide_idx,
                    position=position,
                    filename=path.name,
                )

        subchunks.append(
            SubChunk(
                subchunk_id=len(subchunks),
                content=f" ; End of slide {slide_idx + 1} ** ",
                page=slide_idx,
                position=ChunkPosition(
                    1,
                    1,
                    1,
                    1,
                ),
                file_name=path.name,
            ),
        )

    return subchunks, tables, [], path, True
